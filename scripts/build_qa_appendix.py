from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# Colors
WHITE = RGBColor(255, 255, 255)
TITLE_COLOR = RGBColor(41, 41, 52)
SECTION_COLOR = RGBColor(89, 89, 89)
BODY_COLOR = RGBColor(60, 60, 60)
KEYWORD_COLOR = RGBColor(192, 0, 0)
LINE_COLOR = RGBColor(220, 220, 220)
GAP_COLOR = RGBColor(130, 130, 130)


def add_rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h, text, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT, font="Tahoma"):
    tx_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tx_box


def add_keyword_line(slide, x, y, w, h, keyword, explanation):
    tx_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    run1 = p.add_run()
    run1.text = keyword
    run1.font.name = "Tahoma"
    run1.font.bold = True
    run1.font.color.rgb = KEYWORD_COLOR
    run1.font.size = Pt(13)

    run2 = p.add_run()
    run2.text = " — " + explanation
    run2.font.name = "Tahoma"
    run2.font.bold = False
    run2.font.color.rgb = BODY_COLOR
    run2.font.size = Pt(13)
    return tx_box


def add_slide_shell(prs, section_label, title, note_text):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    add_textbox(
        slide,
        0.5,
        0.15,
        12.3,
        0.2,
        section_label.upper(),
        10,
        SECTION_COLOR,
        bold=True,
    )
    add_textbox(
        slide,
        0.5,
        0.3,
        12.3,
        0.6,
        title,
        22,
        TITLE_COLOR,
        bold=False,
    )
    add_rect(slide, 6.6, 0.9, 0.008, 6.0, LINE_COLOR)

    if note_text:
        add_textbox(slide, 0.5, 6.9, 12.3, 0.4, note_text, 10, SECTION_COLOR, italic=True)
    return slide


def slide_1(prs):
    slide = add_slide_shell(
        prs,
        "Q&A REFERENCE  ·  DATA",
        "Why Mexico Routes Rank Highest — and Why That's Correct",
        "",
    )

    add_keyword_line(
        slide,
        0.5,
        1.0,
        5.8,
        0.78,
        "Fly overlap score 0.92",
        "highest in dataset; Anastrepha ludens, A. obliqua, and Ceratitis capitata all simultaneously present and active in Mexico",
    )
    add_keyword_line(
        slide,
        0.5,
        1.9,
        5.8,
        0.78,
        "LAX cargo volume",
        "Mexico is the #1 source of US fresh fruit imports; LAX receives the highest US–Mexico cargo volume of any port in the dataset",
    )
    add_keyword_line(
        slide,
        0.5,
        2.8,
        5.8,
        0.78,
        "Peak season alignment",
        "August seasonal multiplier is 1.50×; Mexico's summer harvest coincides with peak fly emergence",
    )
    add_keyword_line(
        slide,
        0.5,
        3.7,
        5.8,
        0.78,
        "Diversity toggle built in",
        "click ★ Diverse in the dashboard to force one route per origin country and surface THA, COL, GTM, VNM pathways",
    )

    add_rect(slide, 0.5, 5.8, 5.8, 0.008, LINE_COLOR)
    add_textbox(
        slide,
        0.5,
        5.92,
        5.8,
        0.78,
        "The 10-point gap between Mexico (96.9) and Thailand (86.8) is legitimate. Suppressing it would misrepresent the actual risk.",
        11,
        SECTION_COLOR,
        italic=True,
    )

    # Right table
    add_textbox(slide, 7.0, 1.0, 1.2, 0.3, "Country", 11, SECTION_COLOR, bold=True)
    add_textbox(slide, 8.4, 1.0, 1.2, 0.3, "Fly Overlap", 11, SECTION_COLOR, bold=True)
    add_textbox(slide, 9.8, 1.0, 1.2, 0.3, "Max Risk", 11, SECTION_COLOR, bold=True)
    add_textbox(slide, 11.2, 1.0, 1.6, 0.3, "HIGH Routes", 11, SECTION_COLOR, bold=True)
    add_rect(slide, 7.0, 1.35, 5.8, 0.008, LINE_COLOR)

    rows = [
        ("MEX", "0.92", "96.9", "174"),
        ("THA", "0.88", "86.8", "48"),
        ("GTM", "0.88", "85.4", "33"),
        ("PHL", "0.85", "84.8", "53"),
        ("VNM", "0.82", "83.1", "19"),
        ("COL", "0.82", "82.5", "49"),
        ("SLV", "0.83", "82.2", "12"),
        ("BRA", "0.80", "76.6", "26"),
    ]
    y = 1.5
    for country, overlap, risk, high_routes in rows:
        add_textbox(slide, 7.0, y, 1.2, 0.28, country, 12, BODY_COLOR)
        add_textbox(slide, 8.4, y, 1.2, 0.28, overlap, 12, BODY_COLOR)
        add_textbox(slide, 9.8, y, 1.2, 0.28, risk, 12, BODY_COLOR)
        add_textbox(slide, 11.2, y, 1.2, 0.28, high_routes, 12, BODY_COLOR)
        add_rect(slide, 7.0, y + 0.35, 5.8, 0.008, LINE_COLOR)
        y += 0.42

    add_textbox(slide, 7.0, 5.8, 5.8, 0.3, "Source: data/exports/country_rollup.csv", 10, SECTION_COLOR, italic=True)


def slide_2(prs):
    slide = add_slide_shell(
        prs,
        "Q&A REFERENCE  ·  VALIDATION",
        "How We Validated — and What the Caveats Mean",
        "",
    )

    add_keyword_line(
        slide,
        0.5,
        1.0,
        5.8,
        0.8,
        "In-sample alignment",
        "Pearson r 0.762, Precision@10 100%, HIGH/LOW lift 9.01×; trained and tested on same detection dataset",
    )
    add_keyword_line(
        slide,
        0.5,
        1.9,
        5.8,
        0.8,
        "Temporal holdout",
        "trained on detections ≤2022, tested on ≥2023; Pearson r 0.606, Precision@10 90%; model generalizes to future data it never saw",
    )
    add_keyword_line(
        slide,
        0.5,
        2.8,
        5.8,
        0.8,
        "Leakage sensitivity check",
        "removing detection proximity (10% weight) drops Pearson r to 0.317 but Precision@10 stays 100%; biology and volume carry the signal independently",
    )
    add_keyword_line(
        slide,
        0.5,
        3.7,
        5.8,
        0.8,
        "Surveillance bias disclosed",
        "detections reflect where inspectors trapped, not where flies arrived; HIGH-risk routes with zero detections = candidate trap expansion sites",
    )

    # Metrics quadrants
    add_textbox(slide, 7.0, 1.1, 2.2, 0.45, "0.762", 30, TITLE_COLOR, bold=True)
    add_textbox(slide, 7.0, 1.55, 2.6, 0.25, "Pearson r · in-sample", 11, SECTION_COLOR)

    add_textbox(slide, 10.0, 1.1, 2.0, 0.45, "100%", 30, KEYWORD_COLOR, bold=True)
    add_textbox(slide, 10.0, 1.55, 2.6, 0.25, "Precision@10 · in-sample", 11, SECTION_COLOR)

    add_textbox(slide, 7.0, 2.9, 2.2, 0.45, "0.606", 30, TITLE_COLOR, bold=True)
    add_textbox(slide, 7.0, 3.35, 2.6, 0.25, "Pearson r · holdout 2023+", 11, SECTION_COLOR)

    add_textbox(slide, 10.0, 2.9, 2.0, 0.45, "90%", 30, TITLE_COLOR, bold=True)
    add_textbox(slide, 10.0, 3.35, 2.6, 0.25, "Precision@10 · holdout", 11, SECTION_COLOR)

    add_rect(slide, 7.0, 2.6, 5.8, 0.008, LINE_COLOR)
    add_rect(slide, 9.6, 1.0, 0.008, 3.5, LINE_COLOR)

    add_textbox(slide, 7.0, 4.8, 2.5, 0.3, "Full model", 12, SECTION_COLOR, bold=True)
    add_textbox(slide, 10.0, 4.8, 2.8, 0.3, "No-detection variant", 12, SECTION_COLOR, bold=True)
    add_rect(slide, 7.0, 5.1, 5.8, 0.008, LINE_COLOR)
    add_textbox(slide, 7.0, 5.25, 2.8, 0.25, "Pearson r: 0.762", 12, BODY_COLOR)
    add_textbox(slide, 10.0, 5.25, 2.8, 0.25, "Pearson r: 0.317", 12, BODY_COLOR)
    add_textbox(slide, 7.0, 5.6, 2.8, 0.25, "P@10: 100%", 12, BODY_COLOR)
    add_textbox(slide, 10.0, 5.6, 2.8, 0.25, "P@10: 100%", 12, BODY_COLOR)
    add_textbox(slide, 7.0, 5.95, 2.8, 0.25, "Top-10 overlap", 12, BODY_COLOR)
    add_textbox(slide, 10.0, 5.95, 2.8, 0.25, "50% shared", 12, BODY_COLOR)


def add_formula_block(slide, x, y, w, lines):
    tx_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(1.4))
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.name = "Courier New"
        run.font.size = Pt(12)
        run.font.bold = False
        run.font.italic = False
        run.font.color.rgb = TITLE_COLOR
    return tx_box


def slide_3(prs):
    slide = add_slide_shell(
        prs,
        "Q&A REFERENCE  ·  METHODOLOGY",
        "Risk Score Formula — Every Weight Documented and Tested",
        "",
    )

    add_keyword_line(slide, 0.5, 1.0, 5.8, 0.45, "Structural risk score", "long-term pathway risk; stable month to month")

    add_formula_block(
        slide,
        0.5,
        1.7,
        5.8,
        [
            "structural = 100 × (",
            "  0.35 × fly_overlap",
            "  + 0.25 × volume_score",
            "  + 0.20 × host_fraction",
            "  + 0.10 × route_frequency",
            "  + 0.10 × detection_proximity",
            ") / 0.90",
        ],
    )

    add_rect(slide, 0.5, 3.3, 5.8, 0.008, LINE_COLOR)
    add_keyword_line(slide, 0.5, 3.45, 5.8, 0.45, "Operational surge modifier", "near-term tactical signal")
    add_formula_block(
        slide,
        0.5,
        3.8,
        5.8,
        [
            "surge = 0.85 + 0.30 × (",
            "  0.55 × recent_detection_pressure",
            "  + 0.20 × seasonality",
            "  + 0.25 × port_alert",
            ")",
        ],
    )

    add_rect(slide, 0.5, 5.1, 5.8, 0.008, LINE_COLOR)
    add_keyword_line(slide, 0.5, 5.2, 5.8, 0.45, "Final score", "normalized to prevent ceiling collapse")
    add_formula_block(
        slide,
        0.5,
        5.4,
        5.8,
        ["risk_score = 100 × structural × (surge / 1.15)"],
    )
    add_textbox(slide, 0.5, 6.3, 5.8, 0.25, "Weights sum to 1.0 · enforced in 77 passing tests", 10, SECTION_COLOR, italic=True)

    # Right bars
    bars = [
        ("Fly–host co-occurrence", "35%", 1.93, RGBColor(210, 60, 60), 1.2),
        ("Volume / exposure", "25%", 1.38, RGBColor(80, 80, 100), 1.95),
        ("Host commodity fraction", "20%", 1.10, RGBColor(100, 100, 120), 2.7),
        ("Route frequency", "10%", 0.55, RGBColor(130, 130, 150), 3.45),
        ("Detection proximity", "10%", 0.55, RGBColor(130, 130, 150), 4.2),
    ]
    for label, pct, width, color, y in bars:
        add_textbox(slide, 7.0, y - 0.22, 4.4, 0.2, label, 12, BODY_COLOR)
        add_textbox(slide, 12.2, y - 0.22, 0.6, 0.2, pct, 12, KEYWORD_COLOR, bold=True, align=PP_ALIGN.RIGHT)
        add_rect(slide, 7.0, y, width, 0.22, color)

    add_textbox(slide, 7.0, 5.2, 5.8, 0.25, "Seasonal multiplier: 0.85 (winter) → 1.50 (peak summer)", 11, SECTION_COLOR)
    add_textbox(slide, 7.0, 5.55, 5.8, 0.25, "Southern hemisphere origins: 6-month phase shift applied", 11, SECTION_COLOR)
    add_textbox(slide, 7.0, 5.9, 5.8, 0.25, "Toggle TIER_MODE in phase2_risk_scoring.py for absolute thresholds", 10, SECTION_COLOR, italic=True)


def slide_4(prs):
    slide = add_slide_shell(
        prs,
        "Q&A REFERENCE  ·  DATA SOURCES",
        "Five Sources Integrated — Two Gaps Acknowledged Honestly",
        "",
    )

    add_keyword_line(
        slide,
        0.5,
        1.0,
        5.8,
        0.75,
        "BTS T-100 International Segment",
        "monthly flight-level passenger counts by country × US airport; public domain; transtats.bts.gov",
    )
    add_keyword_line(
        slide,
        0.5,
        1.82,
        5.8,
        0.75,
        "US Census FATUS",
        "monthly cargo by country × commodity × port; HS-6 detail; public domain; usatrade.census.gov",
    )
    add_keyword_line(
        slide,
        0.5,
        2.64,
        5.8,
        0.75,
        "GBIF Tephritidae",
        "species occurrence by country for 13 regulated species; CC-BY 4.0; pipeline degrades gracefully if API unavailable",
    )
    add_keyword_line(
        slide,
        0.5,
        3.46,
        5.8,
        0.75,
        "APHIS detection feature layer",
        "5-year catch records; ground truth for validation; provided at hackathon kickoff",
    )
    add_keyword_line(
        slide,
        0.5,
        4.28,
        5.8,
        0.75,
        "APHIS host commodity list",
        "HS code → host flag mapping; encoded from APHIS PPQ regulatory documents",
    )
    add_rect(slide, 0.5, 5.5, 5.8, 0.008, LINE_COLOR)
    add_keyword_line(
        slide,
        0.5,
        5.62,
        5.8,
        0.45,
        "Express courier — not modeled",
        "no credible public volume data; schema accepts it as transport_mode when data becomes available",
    )
    add_keyword_line(
        slide,
        0.5,
        6.06,
        5.8,
        0.45,
        "USPS mail — not modeled",
        "same reason; acknowledged in dashboard; not fabricated",
    )

    add_textbox(slide, 7.0, 1.0, 1.6, 0.25, "Source", 11, SECTION_COLOR, bold=True)
    add_textbox(slide, 9.0, 2.0, 1.8, 0.25, "Relevance", 11, SECTION_COLOR, bold=True)
    add_textbox(slide, 11.5, 1.0, 1.2, 0.25, "Access", 11, SECTION_COLOR, bold=True)
    add_rect(slide, 7.0, 1.3, 5.8, 0.008, LINE_COLOR)

    rows = [
        ("BTS T-100", "High — route passengers", "Public", False),
        ("FATUS / Census", "High — agricultural imports", "Public", False),
        ("GBIF", "High — species presence", "CC-BY", False),
        ("APHIS detections", "High — ground truth", "Hackathon", False),
        ("Host commodity", "High — host flagging", "Public", False),
        ("Express courier", "Gap — no public data", "—", True),
        ("USPS mail", "Gap — no public data", "—", True),
    ]
    y = 1.45
    for source, rel, access, gap in rows:
        c = GAP_COLOR if gap else BODY_COLOR
        add_textbox(slide, 7.0, y, 1.8, 0.25, source, 12, c, italic=gap)
        add_textbox(slide, 9.0, y, 2.3, 0.25, rel, 12, c, italic=gap)
        add_textbox(slide, 11.5, y, 1.2, 0.25, access, 12, c, italic=gap)
        add_rect(slide, 7.0, y + 0.31, 5.8, 0.008, LINE_COLOR)
        y += 0.48

    add_textbox(
        slide,
        7.0,
        5.8,
        5.8,
        0.4,
        "Live data swap: drop CSVs in data/raw/ → python main.py all → redeploy",
        10,
        SECTION_COLOR,
        italic=True,
    )


def slide_5(prs):
    slide = add_slide_shell(
        prs,
        "Q&A REFERENCE  ·  ARCHITECTURE",
        "Production-Ready Architecture — Already Scaffolded in the Repo",
        "",
    )

    add_keyword_line(
        slide,
        0.5,
        1.0,
        5.8,
        0.72,
        "AWS S3",
        "static dashboard hosting; GeoJSON layers; CSV exports; content-typed uploads; raw source upload opt-in",
    )
    add_keyword_line(
        slide,
        0.5,
        1.78,
        5.8,
        0.72,
        "AWS Step Functions",
        "monthly pipeline orchestration scaffold generated; EventBridge trigger; state machine JSON in infra/aws/",
    )
    add_keyword_line(
        slide,
        0.5,
        2.56,
        5.8,
        0.72,
        "AWS Athena",
        "named queries for port rollup, seasonal analysis, top routes; external tables over S3 CSVs; no ETL required",
    )
    add_keyword_line(
        slide,
        0.5,
        3.34,
        5.8,
        0.72,
        "AWS Lambda",
        "validator function for schema checks; Bedrock narrator for natural language summaries; stubs in repo",
    )
    add_keyword_line(
        slide,
        0.5,
        4.12,
        5.8,
        0.72,
        "ArcGIS Online",
        "4 hosted feature layers published; Dashboard and StoryMap live; overwrite-by-title prevents duplicate sprawl",
    )
    add_keyword_line(
        slide,
        0.5,
        4.90,
        5.8,
        0.72,
        "SageMaker-compatible",
        "model bake-off outputs in CSV format; champion/challenger framework ready for automated promotion",
    )

    # Right flow nodes + lines
    add_textbox(slide, 7.0, 1.1, 1.8, 0.3, "EventBridge", 14, TITLE_COLOR, bold=True)
    add_textbox(slide, 7.0, 1.4, 2.8, 0.25, "1st of month trigger", 11, SECTION_COLOR)
    add_rect(slide, 8.5, 1.7, 0.008, 0.3, LINE_COLOR)

    add_textbox(slide, 7.0, 2.1, 2.0, 0.3, "Step Functions", 14, TITLE_COLOR, bold=True)
    add_textbox(slide, 7.0, 2.4, 2.8, 0.25, "Orchestrate phases 1–4", 11, SECTION_COLOR)
    add_rect(slide, 8.5, 2.7, 0.008, 0.3, LINE_COLOR)

    add_textbox(slide, 7.0, 3.1, 2.0, 0.3, "Athena + S3", 14, TITLE_COLOR, bold=True)
    add_textbox(slide, 7.0, 3.4, 3.0, 0.25, "Query and store scored routes", 11, SECTION_COLOR)

    add_rect(slide, 8.5, 3.7, 0.008, 0.25, LINE_COLOR)
    add_rect(slide, 8.5, 3.95, 1.7, 0.008, LINE_COLOR)  # right split
    add_rect(slide, 7.2, 3.95, 1.3, 0.008, LINE_COLOR)  # left split

    add_textbox(slide, 7.0, 4.1, 1.5, 0.3, "ArcGIS", 13, TITLE_COLOR, bold=True)
    add_textbox(slide, 7.0, 4.38, 1.8, 0.25, "GeoJSON layers", 11, SECTION_COLOR)

    add_textbox(slide, 10.0, 4.1, 2.0, 0.3, "S3 Dashboard", 13, TITLE_COLOR, bold=True)
    add_textbox(slide, 10.0, 4.38, 2.3, 0.25, "Static HTML cockpit", 11, SECTION_COLOR)

    add_textbox(slide, 7.0, 5.5, 5.8, 0.25, "Lambda: schema validation + Bedrock narrative generation", 11, SECTION_COLOR, italic=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, 7.0, 6.1, 5.8, 0.25, "python main.py aws_prod scaffold --bucket <name>", 11, SECTION_COLOR, font="Courier New")


def slide_6(prs):
    slide = add_slide_shell(
        prs,
        "Q&A REFERENCE  ·  ROADMAP",
        "Three Steps from Hackathon Prototype to Operational Tool",
        "",
    )

    add_keyword_line(
        slide,
        0.5,
        1.0,
        5.8,
        0.78,
        "Step 1 — Live data contracts",
        "replace sample CSVs with automated monthly pulls from BTS API and Census FTP; estimated effort: 2 weeks with APHIS data team",
    )
    add_keyword_line(
        slide,
        0.5,
        1.9,
        5.8,
        0.78,
        "Step 2 — IAM and security hardening",
        "production AWS account with least-privilege roles; Secrets Manager; VPC for Lambda; estimated effort: 1 week with cloud team",
    )
    add_keyword_line(
        slide,
        0.5,
        2.8,
        5.8,
        0.78,
        "Step 3 — PPQ entomologist review",
        "validate formula weights against expert judgment; adjust fly overlap scores with live GBIF; champion/challenger promotion policy; estimated effort: 2–4 weeks",
    )
    add_rect(slide, 0.5, 5.2, 5.8, 0.008, LINE_COLOR)
    add_keyword_line(
        slide,
        0.5,
        5.35,
        5.8,
        0.85,
        "What's already done",
        "deterministic pipeline, 77 tests, AWS scaffold generated, ArcGIS layers live, monthly refresh in 30 seconds, documentation complete, methodology published with caveats",
    )

    add_textbox(slide, 7.0, 1.1, 1.1, 0.35, "NOW", 16, KEYWORD_COLOR, bold=True)
    add_textbox(slide, 8.2, 1.1, 2.4, 0.3, "Hackathon prototype", 13, TITLE_COLOR, bold=True)
    add_textbox(slide, 8.2, 1.55, 4.6, 0.25, "Pipeline · Dashboard · Validation · AWS scaffold", 11, SECTION_COLOR)
    add_rect(slide, 7.0, 2.2, 5.8, 0.008, LINE_COLOR)

    add_textbox(slide, 7.0, 2.4, 1.4, 0.35, "MONTH 1", 16, TITLE_COLOR, bold=True)
    add_textbox(slide, 8.5, 2.4, 2.8, 0.3, "Data contracts + security", 13, TITLE_COLOR, bold=True)
    add_textbox(slide, 8.5, 2.85, 4.3, 0.25, "Live BTS/FATUS · IAM hardening · Automated refresh", 11, SECTION_COLOR)
    add_rect(slide, 7.0, 3.5, 5.8, 0.008, LINE_COLOR)

    add_textbox(slide, 7.0, 3.7, 1.6, 0.35, "MONTH 2–3", 16, TITLE_COLOR, bold=True)
    add_textbox(slide, 8.5, 3.7, 2.6, 0.3, "PPQ validation + launch", 13, TITLE_COLOR, bold=True)
    add_textbox(slide, 8.5, 4.15, 4.3, 0.25, "Expert weight review · Champion/challenger · Production SLAs", 11, SECTION_COLOR)
    add_rect(slide, 7.0, 4.8, 5.8, 0.008, LINE_COLOR)

    add_textbox(
        slide,
        7.0,
        5.1,
        5.8,
        0.7,
        "The architecture is production-aligned. The gap is data contracts\nand security hardening — not a rebuild.",
        12,
        SECTION_COLOR,
        italic=True,
    )


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_1(prs)
    slide_2(prs)
    slide_3(prs)
    slide_4(prs)
    slide_5(prs)
    slide_6(prs)

    out_path = "FFED_QA_Appendix.pptx"
    prs.save(out_path)
    print("Saved FFED_QA_Appendix.pptx")

    prs2 = Presentation(out_path)
    for i, slide in enumerate(prs2.slides):
        title = ""
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame and shape.text_frame.text.strip():
                title = shape.text_frame.text.split("\n")[0][:60]
                break
        print(f"[{i+1:02d}] shapes={len(slide.shapes):2d} | {title}")


if __name__ == "__main__":
    build()
